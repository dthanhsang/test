import streamlit as st
import fitz  # PyMuPDF
import zipfile
import io
import os
from pptx import Presentation
from pptx.util import Inches

# --- CẤU HÌNH ---
st.set_page_config(page_title="PDF Converter Pro", page_icon="🚀", layout="centered")

def parse_page_range(range_str, max_pages):
    """Xử lý chọn trang"""
    pages = set()
    if not range_str.strip(): return list(range(max_pages))
    parts = range_str.split(',')
    for part in parts:
        part = part.strip()
        if '-' in part:
            try:
                start, end = map(int, part.split('-'))
                start = max(1, start)
                end = min(max_pages, end)
                for i in range(start, end + 1):
                    pages.add(i - 1)
            except ValueError: continue
        else:
            try:
                p = int(part)
                if 1 <= p <= max_pages:
                    pages.add(p - 1)
            except ValueError: continue
    return sorted(list(pages))

def create_zip_images(doc, pages_to_convert, mat):
    """Tạo file ZIP chứa ảnh"""
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, page_num in enumerate(pages_to_convert):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_data = pix.tobytes("png")
            zf.writestr(f"page_{page_num + 1:03d}.png", img_data)
    zip_buffer.seek(0)
    return zip_buffer

def create_pptx_file(doc, pages_to_convert, mat):
    """Tạo file PowerPoint từ ảnh"""
    prs = Presentation()
    # Layout trắng (số 6)
    blank_slide_layout = prs.slide_layouts[6]
    
    ppt_buffer = io.BytesIO()

    for page_num in pages_to_convert:
        page = doc.load_page(page_num)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_data = pix.tobytes("png")
        
        # Tạo stream ảnh ảo để đưa vào PPT
        image_stream = io.BytesIO(img_data)
        
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Chèn ảnh full chiều rộng slide
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        pic = slide.shapes.add_picture(image_stream, 0, 0, width=slide_width)
        
        # Căn giữa dọc nếu cần
        if pic.height < slide_height:
            pic.top = int((slide_height - pic.height) / 2)

    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

# --- GIAO DIỆN ---
st.title("📄 PDF to Ultra-HD (NGA AN CUC)")
st.markdown("Chuyển đổi PDF sang **Ảnh** hoặc **PowerPoint** chất lượng cao (2K/4K).")

with st.sidebar:
    st.header("⚙️ Cấu hình")
    quality_option = st.selectbox("Chất lượng:", ("Full HD (1080p)", "2K (Siêu nét)", "4K (In ấn)"), index=1)
    zoom_map = {"Full HD (1080p)": 2.0, "2K (Siêu nét)": 3.0, "4K (In ấn)": 4.0}
    zoom_factor = zoom_map[quality_option]
    
    st.divider()
    output_format = st.radio("Định dạng đầu ra:", ["File ZIP (Ảnh rời)", "PowerPoint (.pptx)"])

uploaded_file = st.file_uploader("Tải file PDF lên:", type=["pdf"])

if uploaded_file:
    # Đọc file
    doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
    total_pages = len(doc)
    
    st.info(f"File có **{total_pages} trang**.")
    
    # Chọn trang
    col1, col2 = st.columns([1, 2])
    with col1:
        option = st.radio("Phạm vi:", ["Tất cả", "Tùy chọn"])
    with col2:
        range_input = st.text_input("Nhập trang (VD: 1, 3-5):", value="1" if option=="Tùy chọn" else "", disabled=(option=="Tất cả"))
    
    if st.button("🚀 BẮT ĐẦU XỬ LÝ", type="primary"):
        # Tính toán trang
        pages = list(range(total_pages)) if option == "Tất cả" else parse_page_range(range_input, total_pages)
        
        if not pages:
            st.error("Không có trang nào hợp lệ!")
        else:
            with st.spinner(f"Đang xử lý {len(pages)} trang với độ nét {quality_option}..."):
                mat = fitz.Matrix(zoom_factor, zoom_factor)
                
                # Xử lý theo định dạng chọn
                if output_format == "File ZIP (Ảnh rời)":
                    result_data = create_zip_images(doc, pages, mat)
                    file_ext = "zip"
                    mime_type = "application/zip"
                else:
                    result_data = create_pptx_file(doc, pages, mat)
                    file_ext = "pptx"
                    mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                
                st.success("✅ Xong! NGA AN CUCCC")
                st.download_button(
                    label=f"📥 TẢI VỀ FILE .{file_ext.upper()}",
                    data=result_data,
                    file_name=f"converted_result.{file_ext}",
                    mime=mime_type,
                    use_container_width=True
                )
    

    doc.close()
